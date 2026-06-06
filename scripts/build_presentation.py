"""Generate a polished Open Project 2026 presentation deck (PowerPoint)."""
from __future__ import annotations

import csv
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FIGURES = ROOT / "outputs" / "figures"
METRICS_CSV = ROOT / "outputs" / "evaluation_metrics.csv"
OUT = ROOT / "docs" / "OpenProject2026_Presentation.pptx"
FOOTER_TEXT = "Open Project 2026  |  Agentic AI Dynamic Tariff Optimization"

# ── Theme ─────────────────────────────────────────────────────────────────────
NAVY = RGBColor(13, 27, 58)
NAVY_MID = RGBColor(26, 47, 90)
TEAL = RGBColor(0, 168, 150)
TEAL_LIGHT = RGBColor(224, 247, 244)
ACCENT = RGBColor(255, 183, 3)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(248, 250, 252)
SLATE = RGBColor(71, 85, 105)
DARK_TEXT = RGBColor(30, 41, 59)
MUTED = RGBColor(100, 116, 139)
CARD_BORDER = RGBColor(226, 232, 240)
POSITIVE = RGBColor(5, 150, 105)
NEGATIVE = RGBColor(220, 38, 38)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(1.05)
FOOTER_H = Inches(0.38)
MARGIN = Inches(0.55)
CONTENT_TOP = Inches(1.22)
CONTENT_BOTTOM = Inches(6.95)


def _load_metrics() -> dict[str, float]:
    defaults = {
        "demand_rmse": 0.036,
        "demand_mae": 0.021,
        "demand_r2": 0.942,
        "revenue_gain_pct": -0.67,
        "utilization_before": 0.290,
        "utilization_after": 0.292,
        "off_peak_uplift_pct": 1.11,
        "avg_wait_reduction_pct": 1.35,
        "customer_response_rate_pct": 0.37,
        "pricing_efficiency_baseline": 15.0,
        "pricing_efficiency_dynamic": 14.84,
        "pricing_efficiency_improvement_pct": -1.06,
    }
    if not METRICS_CSV.exists():
        return defaults
    with open(METRICS_CSV, newline="") as f:
        row = next(csv.DictReader(f))
    for key in defaults:
        if key in row and row[key]:
            defaults[key] = float(row[key])
    return defaults


def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape) -> None:
    shape.line.fill.background()


def _font(run, *, size: int, bold: bool = False, color: RGBColor = DARK_TEXT, name: str = "Calibri") -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def _add_rect(slide, left, top, width, height, color: RGBColor, *, line_color: RGBColor | None = None) -> object:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _set_fill(shape, color)
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        _no_line(shape)
    return shape


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = DARK_TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> object:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    _font(p.runs[0], size=size, bold=bold, color=color)
    return box


def _add_rich_textbox(slide, left, top, width, height, lines: list[tuple[str, int, bool, RGBColor]]) -> object:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.space_after = Pt(6)
        if p.runs:
            _font(p.runs[0], size=size, bold=bold, color=color)
        else:
            run = p.add_run()
            run.text = text
            _font(run, size=size, bold=bold, color=color)
    return box


def _slide_base(slide, title: str, *, subtitle: str = "") -> None:
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, OFF_WHITE)
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, HEADER_H, NAVY)
    _add_rect(slide, Inches(0), HEADER_H - Inches(0.06), SLIDE_W, Inches(0.06), TEAL)

    _add_textbox(
        slide, MARGIN, Inches(0.18), Inches(11.5), Inches(0.55),
        title, size=28, bold=True, color=WHITE,
    )
    if subtitle:
        _add_textbox(
            slide, MARGIN, Inches(0.68), Inches(11.5), Inches(0.32),
            subtitle, size=13, color=TEAL_LIGHT,
        )

    _add_rect(slide, Inches(0), SLIDE_H - FOOTER_H, SLIDE_W, FOOTER_H, NAVY_MID)
    _add_textbox(
        slide, MARGIN, SLIDE_H - Inches(0.33), Inches(10), Inches(0.25),
        FOOTER_TEXT, size=9, color=RGBColor(148, 163, 184),
    )


def _add_image(slide, path: Path, left, top, width, *, height=None, caption: str = "") -> None:
    if not path.exists():
        _add_textbox(
            slide, left, top, width, Inches(0.4),
            f"[Figure missing: {path.name}]", size=11, color=MUTED,
        )
        return
    if height is not None:
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        img_bottom = top + height
    else:
        pic = slide.shapes.add_picture(str(path), left, top, width=width)
        img_bottom = top + pic.height
    if caption:
        _add_textbox(slide, left, img_bottom + Inches(0.04), width, Inches(0.32), caption, size=10, color=MUTED, align=PP_ALIGN.CENTER)


def _metric_card(
    slide,
    left,
    top,
    width,
    height,
    label: str,
    value: str,
    *,
    accent: RGBColor = TEAL,
    value_color: RGBColor = NAVY,
    note: str = "",
) -> None:
    card = _add_rect(slide, left, top, width, height, WHITE, line_color=CARD_BORDER)
    _add_rect(slide, left, top, width, Inches(0.07), accent)
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.18), width - Inches(0.3), Inches(0.35), label, size=11, color=MUTED)
    _add_textbox(
        slide, left + Inches(0.15), top + Inches(0.48), width - Inches(0.3), Inches(0.55),
        value, size=26, bold=True, color=value_color, align=PP_ALIGN.LEFT,
    )
    if note:
        _add_textbox(slide, left + Inches(0.15), top + height - Inches(0.42), width - Inches(0.3), Inches(0.35), note, size=9, color=SLATE)


def _bullet_block(slide, left, top, width, height, heading: str, bullets: list[str], *, heading_color: RGBColor = NAVY) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if heading:
        hp = tf.paragraphs[0]
        hp.text = heading
        hp.space_after = Pt(10)
        _font(hp.runs[0], size=15, bold=True, color=heading_color)
    else:
        tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if not heading and i == 0 else tf.add_paragraph()
        p.text = f"•  {bullet}" if bullet else ""
        p.level = 0
        p.space_after = Pt(8)
        if bullet and p.runs:
            _font(p.runs[0], size=13, color=DARK_TEXT)


def _fmt_pct(val: float, *, signed: bool = True) -> str:
    if signed and val > 0:
        return f"+{val:.2f}%"
    return f"{val:.2f}%"


def _build_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    _add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, TEAL)
    _add_rect(slide, Inches(0), Inches(4.85), SLIDE_W, Inches(0.08), TEAL)

    _add_textbox(slide, Inches(1.1), Inches(1.6), Inches(11), Inches(1.2), "Agentic AI Dynamic Tariff Optimization", size=40, bold=True, color=WHITE)
    _add_textbox(
        slide, Inches(1.1), Inches(2.85), Inches(10), Inches(0.6),
        "Intelligent pricing for EV charging networks", size=22, color=TEAL_LIGHT,
    )
    _add_textbox(
        slide, Inches(1.1), Inches(3.55), Inches(10), Inches(0.5),
        "Open Project 2026", size=18, bold=True, color=ACCENT,
    )
    _add_textbox(
        slide, Inches(1.1), Inches(5.2), Inches(10), Inches(0.4),
        "UrbanEV (ST-EVCDP)  +  ACN Workplace Charging Datasets", size=14, color=RGBColor(148, 163, 184),
    )
    _add_textbox(
        slide, Inches(1.1), Inches(5.65), Inches(10), Inches(0.4),
        "Demand forecasting  ·  Dynamic tariffs  ·  Closed-loop monitoring", size=13, color=RGBColor(148, 163, 184),
    )


def _build_problem_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_base(slide, "The Challenge", subtitle="Why dynamic tariffs matter for EV charging operators")

    col_w = Inches(5.85)
    gap = Inches(0.45)
    left_x = MARGIN
    right_x = MARGIN + col_w + gap

    _bullet_block(
        slide, left_x, CONTENT_TOP, col_w, Inches(2.8),
        "Market context",
        [
            "EV adoption is shifting load profiles — workplace and urban hubs face recurring congestion peaks.",
            "Flat pricing cannot signal scarcity or incentivize off-peak charging.",
            "Operators must balance revenue, grid stability, and driver experience.",
        ],
    )
    _bullet_block(
        slide, left_x, Inches(3.55), col_w, Inches(2.8),
        "Our approach",
        [
            "Agentic AI pipeline: forecast demand, set bounded dynamic tariffs, evaluate KPIs.",
            "Dual-dataset validation — UrbanEV (Shenzhen grids) + ACN (Caltech/JPL workplace).",
            "Consumer-acceptable pricing guardrails: max ±25% vs baseline, 10% max step change.",
        ],
    )

    panel = _add_rect(slide, right_x, CONTENT_TOP, col_w, Inches(5.5), WHITE, line_color=CARD_BORDER)
    _add_textbox(slide, right_x + Inches(0.25), CONTENT_TOP + Inches(0.2), col_w - Inches(0.5), Inches(0.35), "Target outcomes", size=15, bold=True, color=NAVY)
    outcomes = [
        ("Predict next-hour utilization", "HistGradientBoosting on UrbanEV hourly panel"),
        ("Shape demand with smart tariffs", "Surge at 75% util · Discount below 35%"),
        ("Measure operational impact", "Revenue, off-peak uplift, wait proxy, elasticity"),
        ("Learn from feedback", "Monitoring agent tunes thresholds over time"),
    ]
    y = CONTENT_TOP + Inches(0.7)
    for title, desc in outcomes:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, right_x + Inches(0.35), y + Inches(0.08), Inches(0.14), Inches(0.14))
        _set_fill(dot, TEAL)
        _no_line(dot)
        _add_textbox(slide, right_x + Inches(0.6), y, col_w - Inches(0.85), Inches(0.3), title, size=13, bold=True, color=NAVY)
        _add_textbox(slide, right_x + Inches(0.6), y + Inches(0.32), col_w - Inches(0.85), Inches(0.45), desc, size=11, color=SLATE)
        y += Inches(1.15)


def _build_data_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_base(slide, "Data Landscape & Preprocessing", subtitle="Two complementary datasets powering demand and revenue simulation")

    cards = [
        ("ACN Workplace", "14,947 sessions", "Caltech / JPL\nconnectionTime, kWh, duration"),
        ("UrbanEV Grids", "247 spatial grids", "Shenzhen ST-EVCDP\n5-min → hourly aggregates"),
        ("Engineered Features", "Temporal + spatial", "Utilization, lags, congestion proxy"),
        ("Quality Controls", "Validated pipeline", "Drop zero-kWh, invalid durations"),
    ]
    card_w = Inches(2.85)
    card_h = Inches(1.55)
    x0 = MARGIN
    for i, (title, stat, detail) in enumerate(cards):
        x = x0 + i * (card_w + Inches(0.22))
        _add_rect(slide, x, CONTENT_TOP, card_w, card_h, WHITE, line_color=CARD_BORDER)
        _add_rect(slide, x, CONTENT_TOP, card_w, Inches(0.06), TEAL if i % 2 == 0 else NAVY_MID)
        _add_textbox(slide, x + Inches(0.15), CONTENT_TOP + Inches(0.18), card_w - Inches(0.3), Inches(0.3), title, size=13, bold=True, color=NAVY)
        _add_textbox(slide, x + Inches(0.15), CONTENT_TOP + Inches(0.48), card_w - Inches(0.3), Inches(0.35), stat, size=18, bold=True, color=TEAL)
        _add_textbox(slide, x + Inches(0.15), CONTENT_TOP + Inches(0.88), card_w - Inches(0.3), Inches(0.55), detail, size=10, color=SLATE)

    img_top = Inches(3.05)
    img_w = Inches(5.9)
    _add_image(slide, FIGURES / "acn_weekday_weekend.png", MARGIN, img_top, img_w, caption="ACN: weekday vs weekend session patterns")
    _add_image(
        slide, FIGURES / "urbanev_price_vs_util.png", MARGIN + img_w + Inches(0.35), img_top, img_w,
        caption="UrbanEV: price vs utilization relationship across grids",
    )


def _build_eda_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_base(slide, "EDA: Demand Behavior Insights", subtitle="Peak timing and utilization patterns inform tariff thresholds")

    insight_w = Inches(3.5)
    _add_rect(slide, MARGIN, CONTENT_TOP, insight_w, Inches(5.45), WHITE, line_color=CARD_BORDER)
    _add_textbox(slide, MARGIN + Inches(0.2), CONTENT_TOP + Inches(0.2), insight_w - Inches(0.4), Inches(0.35), "Key findings", size=15, bold=True, color=NAVY)

    findings = [
        ("ACN workplace peaks", "9–11 AM and 1–3 PM align with arrival and mid-day top-up behavior."),
        ("UrbanEV evening load", "Weekday congestion builds 5–9 PM across Shenzhen grids."),
        ("Threshold calibration", "Surge at 75% and discount at 35% — earlier than default 80/30 to act before queues form."),
        ("Revenue baseline", "INR 15/kWh reference tariff for simulation and guardrails."),
    ]
    y = CONTENT_TOP + Inches(0.65)
    for title, body in findings:
        bar = _add_rect(slide, MARGIN + Inches(0.2), y, Inches(0.05), Inches(0.85), TEAL)
        _add_textbox(slide, MARGIN + Inches(0.35), y, insight_w - Inches(0.55), Inches(0.28), title, size=12, bold=True, color=NAVY)
        _add_textbox(slide, MARGIN + Inches(0.35), y + Inches(0.3), insight_w - Inches(0.55), Inches(0.55), body, size=11, color=SLATE)
        y += Inches(1.15)

    img_left = MARGIN + insight_w + Inches(0.35)
    img_w = Inches(4.15)
    _add_image(slide, FIGURES / "acn_hourly_demand.png", img_left, CONTENT_TOP, img_w, caption="ACN hourly demand profile")
    _add_image(slide, FIGURES / "urbanev_util_heatmap.png", img_left + img_w + Inches(0.25), CONTENT_TOP, img_w, caption="UrbanEV utilization heatmap (hour × day)")


def _build_architecture_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_base(slide, "System Architecture", subtitle="Agentic pipeline with forecast → price → monitor → learn loop")

    box_w = Inches(2.35)
    box_h = Inches(1.35)
    y = Inches(2.0)
    stages = [
        ("1. Preprocess", "ACN sessions\nUrbanEV panel\nHourly aggregates", NAVY),
        ("2. Demand Model", "HistGradient\nBoosting\nNext-hour util", NAVY_MID),
        ("3. Tariff Agent", "Util → price map\nSmooth bounds\nACN revenue sim", TEAL),
        ("4. Monitor", "6 KPI metrics\nFeedback loop\nThreshold tune", NAVY_MID),
    ]
    x = Inches(0.75)
    for i, (title, body, color) in enumerate(stages):
        shape = _add_rect(slide, x, y, box_w, box_h, color)
        _add_textbox(slide, x + Inches(0.12), y + Inches(0.12), box_w - Inches(0.24), Inches(0.3), title, size=13, bold=True, color=WHITE)
        _add_textbox(slide, x + Inches(0.12), y + Inches(0.45), box_w - Inches(0.24), Inches(0.8), body, size=10, color=RGBColor(226, 232, 240))
        if i < len(stages) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + box_w + Inches(0.08), y + Inches(0.5), Inches(0.45), Inches(0.35))
            _set_fill(arrow, ACCENT)
            _no_line(arrow)
        x += box_w + Inches(0.55)

    # Feedback loop arc (simplified as labeled band)
    loop = _add_rect(slide, Inches(0.75), Inches(3.65), Inches(11.8), Inches(0.55), TEAL_LIGHT, line_color=TEAL)
    _add_textbox(
        slide, Inches(0.95), Inches(3.72), Inches(11.4), Inches(0.4),
        "↺  Monitoring agent feeds threshold adjustments back to Tariff Agent (surge/discount offsets)",
        size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
    )

    # Lower detail panels
    panel_y = Inches(4.55)
    panel_w = Inches(5.85)
    _add_rect(slide, MARGIN, panel_y, panel_w, Inches(2.15), WHITE, line_color=CARD_BORDER)
    _add_textbox(slide, MARGIN + Inches(0.2), panel_y + Inches(0.15), panel_w - Inches(0.4), Inches(0.3), "Demand + Tariff Agent (merged)", size=14, bold=True, color=NAVY)
    agent_lines = [
        "Predict utilization on UrbanEV test window (last 20% chronologically)",
        "Map util → tariff: surge ≤ +25%, discount ≤ −15% vs INR 15/kWh baseline",
        "Apply 10% max step smoothing to avoid price shocks",
        "Simulate ACN revenue shift using price elasticity (−0.35)",
    ]
    _bullet_block(slide, MARGIN + Inches(0.15), panel_y + Inches(0.5), panel_w - Inches(0.3), Inches(1.55), "", agent_lines)

    right_x = MARGIN + panel_w + Inches(0.45)
    _add_rect(slide, right_x, panel_y, panel_w, Inches(2.15), WHITE, line_color=CARD_BORDER)
    _add_textbox(slide, right_x + Inches(0.2), panel_y + Inches(0.15), panel_w - Inches(0.4), Inches(0.3), "Monitoring & Learning Agent", size=14, bold=True, color=NAVY)
    kpi_lines = [
        "Revenue Gain % (ACN) · Off-Peak Uplift (UrbanEV)",
        "Wait Reduction proxy · Customer Response Rate",
        "Pricing Efficiency (INR/kWh) · Utilization before/after",
        "Outputs: evaluation_metrics.csv, monitoring_feedback.json",
    ]
    _bullet_block(slide, right_x + Inches(0.15), panel_y + Inches(0.5), panel_w - Inches(0.3), Inches(1.55), "", kpi_lines)


def _build_model_slide(prs: Presentation, metrics: dict[str, float]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_base(slide, "Demand Prediction & Dynamic Tariff", subtitle="High-accuracy forecasting drives bounded, consumer-friendly pricing")

    rmse = metrics["demand_rmse"]
    mae = metrics["demand_mae"]
    r2 = metrics["demand_r2"]

    cards = [
        ("RMSE", f"{rmse:.3f}", "Next-hour utilization error"),
        ("MAE", f"{mae:.3f}", "Mean absolute deviation"),
        ("R²", f"{r2:.3f}", "Explained variance"),
    ]
    card_w = Inches(2.0)
    for i, (label, val, note) in enumerate(cards):
        _metric_card(slide, MARGIN + i * (card_w + Inches(0.2)), CONTENT_TOP, card_w, Inches(1.25), label, val, note=note)

    # Tariff policy callout
    policy_x = MARGIN + Inches(6.5)
    policy_w = Inches(5.9)
    _add_rect(slide, policy_x, CONTENT_TOP, policy_w, Inches(1.25), NAVY)
    _add_textbox(slide, policy_x + Inches(0.2), CONTENT_TOP + Inches(0.15), policy_w - Inches(0.4), Inches(0.28), "Tariff policy", size=12, bold=True, color=TEAL_LIGHT)
    _add_textbox(
        slide, policy_x + Inches(0.2), CONTENT_TOP + Inches(0.45), policy_w - Inches(0.4), Inches(0.7),
        "Surge ≥ 75% util  ·  Discount ≤ 35% util  ·  Max ±25% vs baseline  ·  10% step cap",
        size=13, color=WHITE,
    )

    img_top = Inches(2.65)
    _add_image(slide, FIGURES / "demand_metrics.png", MARGIN, img_top, Inches(4.2), caption="Model fit on UrbanEV hold-out set")
    _add_image(slide, FIGURES / "feature_importance.png", MARGIN + Inches(4.45), img_top, Inches(4.35), caption="Top predictive features")
    _add_image(slide, FIGURES / "tariff_timeline.png", MARGIN + Inches(8.95), img_top, Inches(3.85), caption="Recommended tariff vs utilization over time")


def _build_results_slide(prs: Presentation, metrics: dict[str, float]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_base(slide, "Results & KPI Dashboard", subtitle="Operational gains with a modest, intentional revenue trade-off")

    rev_color = NEGATIVE if metrics["revenue_gain_pct"] < 0 else POSITIVE
    cards = [
        ("Revenue Gain", _fmt_pct(metrics["revenue_gain_pct"]), "ACN simulation", rev_color),
        ("Off-Peak Uplift", _fmt_pct(metrics["off_peak_uplift_pct"]), "UrbanEV low-util slots", POSITIVE),
        ("Wait Reduction", _fmt_pct(metrics["avg_wait_reduction_pct"]), "Congestion proxy", POSITIVE),
        ("Utilization", f"{metrics['utilization_before']:.1%} → {metrics['utilization_after']:.1%}", "Network-wide", TEAL),
        ("Customer Response", _fmt_pct(metrics["customer_response_rate_pct"]), "Elasticity model", TEAL),
        ("Pricing Efficiency", f"₹{metrics['pricing_efficiency_baseline']:.2f} → ₹{metrics['pricing_efficiency_dynamic']:.2f}", "INR per kWh", SLATE),
    ]
    card_w = Inches(1.95)
    card_h = Inches(1.35)
    for i, (label, val, note, accent) in enumerate(cards):
        row, col = divmod(i, 3)
        x = MARGIN + col * (card_w + Inches(0.22))
        y = CONTENT_TOP + row * (card_h + Inches(0.18))
        val_color = accent if accent in (POSITIVE, NEGATIVE) else NAVY
        _metric_card(slide, x, y, card_w, card_h, label, val, accent=accent if accent not in (POSITIVE, NEGATIVE) else TEAL, value_color=val_color, note=note)

    _add_image(slide, FIGURES / "monitoring_kpis.png", MARGIN, Inches(4.15), Inches(7.2), caption="Monitoring agent KPI summary")
    _add_image(slide, FIGURES / "urbanev_util_distribution.png", MARGIN + Inches(7.45), Inches(4.15), Inches(4.85), caption="Utilization distribution after dynamic pricing")

    _add_rect(slide, MARGIN + Inches(7.45), Inches(6.35), Inches(4.85), Inches(0.55), TEAL_LIGHT, line_color=TEAL)
    _add_textbox(
        slide, MARGIN + Inches(7.6), Inches(6.42), Inches(4.55), Inches(0.4),
        "−0.67% revenue traded for +1.11% off-peak uplift and +1.35% wait reduction",
        size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
    )


def _build_implications_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_base(slide, "Business Impact & Roadmap", subtitle="Actionable insights for operators, policymakers, and future research")

    col_w = Inches(5.85)
    gap = Inches(0.45)
    left_x = MARGIN
    right_x = MARGIN + col_w + gap

    _add_rect(slide, left_x, CONTENT_TOP, col_w, Inches(5.45), WHITE, line_color=CARD_BORDER)
    _add_textbox(slide, left_x + Inches(0.2), CONTENT_TOP + Inches(0.15), col_w - Inches(0.4), Inches(0.35), "Business & policy implications", size=15, bold=True, color=NAVY)
    implications = [
        "Dynamic tariffs flatten peaks without extreme price shocks — feasible for consumer acceptance.",
        "Off-peak discounts align charging with grid renewable surplus and under-utilized capacity.",
        "Operators face a clear trade-off: modest revenue dip for congestion relief and smoother load.",
        "Cross-geography insight: ACN (US workplace) patterns validate revenue logic; UrbanEV (Shenzhen) drives operations KPIs.",
    ]
    y = CONTENT_TOP + Inches(0.6)
    for text in implications:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left_x + Inches(0.3), y + Inches(0.1), Inches(0.12), Inches(0.12))
        _set_fill(dot, TEAL)
        _no_line(dot)
        _add_textbox(slide, left_x + Inches(0.52), y, col_w - Inches(0.75), Inches(0.85), text, size=12, color=DARK_TEXT)
        y += Inches(1.05)

    _add_rect(slide, right_x, CONTENT_TOP, col_w, Inches(5.45), WHITE, line_color=CARD_BORDER)
    _add_textbox(slide, right_x + Inches(0.2), CONTENT_TOP + Inches(0.15), col_w - Inches(0.4), Inches(0.35), "Limitations & future work", size=15, bold=True, color=NAVY)
    future = [
        ("Price elasticity assumed", "−0.35 short-run estimate; not causally identified from ACN alone."),
        ("Wait time proxy", "Occupancy/capacity ratio — not observed queue lengths."),
        ("Geographic transfer", "US workplace + China urban grids require careful generalization."),
        ("Next steps", "RL-based policies, A/B field tests, CO₂-aware grid signals."),
    ]
    y = CONTENT_TOP + Inches(0.6)
    for title, body in future:
        _add_textbox(slide, right_x + Inches(0.25), y, col_w - Inches(0.5), Inches(0.28), title, size=12, bold=True, color=NAVY_MID)
        _add_textbox(slide, right_x + Inches(0.25), y + Inches(0.3), col_w - Inches(0.5), Inches(0.55), body, size=11, color=SLATE)
        y += Inches(1.15)

    _add_textbox(
        slide, MARGIN, Inches(6.55), Inches(12), Inches(0.35),
        "Reproducible pipeline:  python main.py  →  python scripts/build_presentation.py",
        size=11, color=MUTED, align=PP_ALIGN.CENTER,
    )


def _build_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), TEAL)

    _add_textbox(slide, Inches(1.0), Inches(2.4), Inches(11), Inches(0.9), "Thank You", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(
        slide, Inches(1.0), Inches(3.45), Inches(11), Inches(0.5),
        "Agentic AI Dynamic Tariff Optimization — Open Project 2026", size=18, color=TEAL_LIGHT, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, Inches(1.0), Inches(4.2), Inches(11), Inches(0.45),
        "RMSE 0.036  ·  R² 0.942  ·  Off-peak uplift +1.11%  ·  Wait reduction +1.35%", size=14, color=ACCENT, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide, Inches(1.0), Inches(5.5), Inches(11), Inches(0.4),
        "Questions & discussion", size=16, color=RGBColor(148, 163, 184), align=PP_ALIGN.CENTER,
    )


def build() -> Path:
    metrics = _load_metrics()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _build_title_slide(prs)
    _build_problem_slide(prs)          # Content 1
    _build_data_slide(prs)             # Content 2
    _build_eda_slide(prs)              # Content 3
    _build_architecture_slide(prs)     # Content 4
    _build_model_slide(prs, metrics)   # Content 5
    _build_results_slide(prs, metrics) # Content 6
    _build_implications_slide(prs)     # Content 7
    _build_closing_slide(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    size_kb = path.stat().st_size / 1024
    print(f"Saved: {path} ({size_kb:.0f} KB)")
